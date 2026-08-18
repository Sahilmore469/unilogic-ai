"""
Unilogic AI - Presentation Generator
Creates a professional, dark-themed 8-slide PowerPoint presentation (.pptx) detailing the Unilogic AI solution.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Colors
    BG_DARK = RGBColor(15, 23, 42)       # #0F172A
    CARD_BG = RGBColor(30, 41, 59)       # #1E293B
    TEXT_WHITE = RGBColor(248, 250, 252) # #F8FAFC
    TEXT_MUTED = RGBColor(148, 163, 184)# #94A3B8
    BLUE_ACCENT = RGBColor(59, 130, 246)# #3B82F6
    CYAN_ACCENT = RGBColor(6, 182, 212)  # #06B6D4
    GOLD_ACCENT = RGBColor(245, 158, 11) # #F59E0B
    EMERALD_ACCENT = RGBColor(16, 185, 129)# #10B981

    def apply_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return slide

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BLUE_ACCENT):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    apply_bg(slide1)
    
    add_card(slide1, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5), bg_color=CARD_BG, border_color=BLUE_ACCENT)
    
    txBox = slide1.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.333), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Unilogic AI"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    
    p2 = tf.add_paragraph()
    p2.text = "AI-Powered Product Intelligence for Industrial Commerce"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(14)
    
    p3 = tf.add_paragraph()
    p3.text = "Transforming minimal, cryptic supplier strings into rich 252-column commerce-ready product catalogs."
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_MUTED
    p3.space_before = Pt(16)

    # SLIDE 2: Problem & Challenge
    slide2 = prs.slides.add_slide(blank_layout)
    apply_bg(slide2)
    
    # Title
    txBox = slide2.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "The Challenge: Messy Raw Catalog Data"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    # 3 Cards
    cards_data = [
        ("Cryptic & Abbreviated", "Descriptions like '3/8 CPLG BRS 150#' are fragmented, missing key attributes and unsearchable by end buyers.", GOLD_ACCENT),
        ("Supplier & Brand Noise", "Same manufacturer under 6 spellings ('Appliance Dealers Cooperative (APPDE)'). Field values filled with '-- Unbranded --'.", CYAN_ACCENT),
        ("Strict Formatting Rules", "5 description variants required (Invoice Desc ≤40 UPPERCASE), UOM spacing rules ('24 in'), and 63 inch decimal conversions.", EMERALD_ACCENT)
    ]
    
    for i, (title, desc, color) in enumerate(cards_data):
        left = Inches(1.0 + i * 3.9)
        add_card(slide2, left, Inches(2.0), Inches(3.6), Inches(4.5), border_color=color)
        
        tb = slide2.shapes.add_textbox(left + Inches(0.2), Inches(2.3), Inches(3.2), Inches(3.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(15)
        p2.font.color.rgb = TEXT_MUTED
        p2.space_before = Pt(14)

    # SLIDE 3: System Architecture & 7 AI Agents
    slide3 = prs.slides.add_slide(blank_layout)
    apply_bg(slide3)
    
    txBox = slide3.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture: 7 Sequential AI Agents"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    agents = [
        "1. Ingestion & De-duplication (Filter '-- Unbranded --' placeholders)",
        "2. Entity Resolution (Canonical Rheem Manufacturing & FRIGIDAIRE®)",
        "3. Taxonomy Classification (Dept > Class > Fine > UNSPSC > Classpath)",
        "4. LOV & UOM Normalizer (500+ UOMs: '24 in', '120 V', '15 A')",
        "5. Decimal-to-Fraction Engine (63 fractions: 50.25 in -> 50-1/4 in)",
        "6. Content Description Builder (Invoice ≤40 UPPERCASE, Mobile 60-80)",
        "7. Autonomous Audit Agent (Confidence scoring & HITL review queue)"
    ]

    add_card(slide3, Inches(1.0), Inches(1.8), Inches(11.333), Inches(5.0), border_color=BLUE_ACCENT)
    tb = slide3.shapes.add_textbox(Inches(1.3), Inches(2.1), Inches(10.7), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, ag in enumerate(agents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ag
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(10)

    # SLIDE 4: Entity Resolution & Master UOM Engine
    slide4 = prs.slides.add_slide(blank_layout)
    apply_bg(slide4)
    
    txBox = slide4.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Innovation: Entity Resolution & UOM Standards"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    add_card(slide4, Inches(1.0), Inches(1.8), Inches(5.4), Inches(5.0), border_color=CYAN_ACCENT)
    tb1 = slide4.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(5.0), Inches(4.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "Brand Entity Resolution"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT
    
    p2 = tf1.add_paragraph()
    p2.text = "• Fuzzy matching against 27,000+ UniCat canonical entries.\n• Removes supplier noise & codes (e.g. 'APPDE').\n• Preserves legal symbols (® and ™) e.g., 'FRIGIDAIRE®'."
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(12)

    add_card(slide4, Inches(6.9), Inches(1.8), Inches(5.4), Inches(5.0), border_color=EMERALD_ACCENT)
    tb2 = slide4.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.0), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Master UOM Standardizer"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT
    
    p2 = tf2.add_paragraph()
    p2.text = "• 500+ approved UOM abbreviations mapped.\n• Strict Single-Space Rule: '24 in' (not '24in'), '120 V' (not '120V').\n• Compound dimension parsing: '24 in W x 24-1/4 in D'."
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(12)

    # SLIDE 5: Decimal to Fraction & Description Rules
    slide5 = prs.slides.add_slide(blank_layout)
    apply_bg(slide5)
    
    txBox = slide5.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Trade Fractions & 5 Description Formats"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    add_card(slide5, Inches(1.0), Inches(1.8), Inches(11.333), Inches(5.0), border_color=GOLD_ACCENT)
    tb = slide5.shapes.add_textbox(Inches(1.3), Inches(2.1), Inches(10.7), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    items = [
        ("63 Inch Fraction Conversion:", "Manufacturers publish decimals (50.25); trade buyers search fractions. Automatic exact 1/64 to 63/64 conversion ('50-1/4 in')."),
        ("INVOICE_DESC (≤40 Chars, UPPERCASE):", "Strict character limit & UPPERCASE validation: 'DISHWASHER LEG 5 SST 120V 15A 50-1/4IN'."),
        ("MOBILE_DESC (60-80 Chars):", "Target length formula: '[Manufacturer] [Brand], [Product Type], [Series], [MPN]'."),
        ("SHORT_DESC / Product Title:", "[Brand] [Series] [MPN] [Product Name] With [Feature], [Key Attributes]."),
        ("LONG_DESC & RETAIL_DESC:", "Full technical summary with standard UOMs + consumer marketing summary.")
    ]
    
    for i, (heading, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{heading} "
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = GOLD_ACCENT
        p.space_before = Pt(8)
        
        run = p.add_run()
        run.text = text
        run.font.bold = False
        run.font.color.rgb = TEXT_MUTED

    # SLIDE 6: Ground Truth Benchmark Scorecard
    slide6 = prs.slides.add_slide(blank_layout)
    apply_bg(slide6)
    
    txBox = slide6.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Ground Truth Benchmark Scoring"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    scores = [
        ("100.0%", "Invoice Desc Compliance", "≤40 Chars & 100% UPPERCASE", EMERALD_ACCENT),
        ("100.0%", "Mobile Desc Compliance", "60–80 Chars Target Window", EMERALD_ACCENT),
        ("68.95%", "Overall Alignment Score", "252 Delivery Format Fields", BLUE_ACCENT),
        ("65.67%", "Fuzzy Match Accuracy", "Semantic Entity Alignment", CYAN_ACCENT)
    ]
    
    for i, (val, title, detail, color) in enumerate(scores):
        left = Inches(1.0 + (i % 2) * 5.8)
        top = Inches(1.8 + (i // 2) * 2.6)
        add_card(slide6, left, top, Inches(5.4), Inches(2.3), border_color=color)
        
        tb = slide6.shapes.add_textbox(left + Inches(0.3), top + Inches(0.2), Inches(4.8), Inches(1.9))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_WHITE
        
        p3 = tf.add_paragraph()
        p3.text = detail
        p3.font.size = Pt(14)
        p3.font.color.rgb = TEXT_MUTED

    # SLIDE 7: Web Application Studio
    slide7 = prs.slides.add_slide(blank_layout)
    apply_bg(slide7)
    
    txBox = slide7.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Full-Stack Glassmorphic Web App"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    add_card(slide7, Inches(1.0), Inches(1.8), Inches(11.333), Inches(5.0), border_color=BLUE_ACCENT)
    tb = slide7.shapes.add_textbox(Inches(1.3), Inches(2.1), Inches(10.7), Inches(4.4))
    tf = tb.text_frame
    tf.word_wrap = True
    
    features = [
        "• Pipeline Studio: Real-time 7-agent step visualizer with live character length meters.",
        "• Batch Catalog Engine: Process 1,000 catalog rows with search, filter, and 252-column CSV download.",
        "• Ground Truth Benchmark View: Real-time scorecards and rule compliance audit.",
        "• Human-In-The-Loop (HITL) Studio: Review queue for data stewards to inspect & override flagged rows.",
        "• Master Guidelines Viewer: Interactive content formula & UOM reference cheat sheet."
    ]
    
    for i, feat in enumerate(features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = feat
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(14)

    # SLIDE 8: Enterprise Scalability & Conclusion
    slide8 = prs.slides.add_slide(blank_layout)
    apply_bg(slide8)
    
    txBox = slide8.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Enterprise Value & Scalability"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    add_card(slide8, Inches(1.5), Inches(1.8), Inches(10.333), Inches(4.8), border_color=EMERALD_ACCENT)
    tb = slide8.shapes.add_textbox(Inches(1.8), Inches(2.1), Inches(9.7), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    concl = [
        "1. Instant Scalability: Transform millions of raw distributor catalog rows in minutes.",
        "2. Zero Invented Facts: 100% constrained within Unilog Master UOM & LOV vocabularies.",
        "3. REST API Ready: Plug-and-play FastAPI integration for PIM, ERP, and e-commerce platforms.",
        "4. Measurable Accuracy: 100% constraint pass rate backed by labeled ground truth evaluation."
    ]
    
    for i, c in enumerate(concl):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = c
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.space_before = Pt(16)

    # Save presentation
    output_path = r'C:\Users\Sahil more\.gemini\antigravity\scratch\product-intelligence-engine\Unilogic_AI_Product_Intelligence.pptx'
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
